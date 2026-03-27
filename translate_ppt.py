from pptx import Presentation
from deep_translator import GoogleTranslator
import copy

# ===== SETTINGS =====
input_file = "input.pptx"
output_file = "bilingual_output.pptx"

skip_slides = 50   # Already completed slides (25 EN + 25 TA)
target_lang = "ta"

translator = GoogleTranslator(source="auto", target=target_lang)

prs = Presentation(input_file)
new_prs = Presentation()

def copy_slide(source_slide, target_prs):
    layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(layout)

    for shape in source_slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide


def translate_slide(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    try:
                        translated = translator.translate(text)
                        paragraph.text = translated
                    except:
                        print("Error translating:", text)


for i, slide in enumerate(prs.slides):
    print(f"Processing Slide {i+1}")

    # Copy original
    new_slide_en = copy_slide(slide, new_prs)

    if i >= skip_slides:
        # Create Tamil version
        new_slide_ta = copy_slide(slide, new_prs)
        translate_slide(new_slide_ta)

new_prs.save(output_file)

print("✅ Done! Check bilingual_output.pptx")